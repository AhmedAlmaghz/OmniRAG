# qa_deck.py — فحص جودة برمجي للعرض: تجاوز النص، التداخل، النصوص الناقصة
# overflow: estimated wrapped-text height vs. shape height (Arabic ~char-width heuristic)
# overlap: bounding-box intersection between text-bearing shapes (excluding intentional fills)
from pptx import Presentation
from pptx.util import Emu
import re, unicodedata

EMU_IN = 914400
prs = Presentation("Evolution-of-AI-Agents-AR.pptx")
SLIDE_W, SLIDE_H = prs.slide_width / EMU_IN, prs.slide_height / EMU_IN

# Arabic/Latin mixed width heuristic: avg glyph width factor of font size (inches per char at 1pt)
W_FACTOR = 0.0075   # ~0.54pt avg advance per 1pt font (Arabic is narrow, Segoe UI balanced)

def txt_shapes(shapes):
    for sh in shapes:
        if sh.shape_type == 6:
            yield from txt_shapes(sh.shapes)
        elif sh.has_text_frame and sh.text_frame.text.strip():
            yield sh

def est_height(sh):
    """Estimate rendered text height in inches for a wrapped text box."""
    w_in = sh.width / EMU_IN
    h_in = sh.height / EMU_IN
    total_lines = 0.0
    max_sz = 0
    for para in sh.text_frame.paragraphs:
        text = "".join(r.text for r in para.runs)
        if not text:
            continue
        sz = None
        for r in para.runs:
            if r.font.size:
                sz = max(sz or 0, r.font.size.pt)
        sz = sz or 12
        max_sz = max(max_sz, sz)
        # chars per line: width / (char_width = sz * W_FACTOR)
        cpl = max(1, int(w_in / (sz * W_FACTOR)))
        total_lines += max(1, -(-len(text) // cpl))
    if total_lines == 0:
        return 0, 0
    # line height ~ fontSize * 1.35 (incl. lineSpacing set explicitly)
    line_h = max_sz * 1.35 / 72
    return total_lines * line_h, max_sz

issues = []
for si, slide in enumerate(prs.slides, 1):
    shapes = list(txt_shapes(slide.shapes))
    for sh in shapes:
        # off-slide check
        x, y = sh.left / EMU_IN, sh.top / EMU_IN
        w, h = sh.width / EMU_IN, sh.height / EMU_IN
        if x < -0.05 or y < -0.05 or x + w > SLIDE_W + 0.05 or y + h > SLIDE_H + 0.05:
            issues.append(f"S{si} OFF-SLIDE: '{sh.text_frame.text[:25]}' at ({x:.2f},{y:.2f},{w:.2f},{h:.2f})")
        eh, sz = est_height(sh)
        if eh > 0 and eh > h * 1.18:
            issues.append(f"S{si} OVERFLOW? '{sh.text_frame.text[:30]}' est {eh:.2f}in > box {h:.2f}in (fs~{sz})")
    # overlap check among text shapes
    for i in range(len(shapes)):
        for j in range(i + 1, len(shapes)):
            a, b = shapes[i], shapes[j]
            ax, ay, aw, ah = a.left/EMU_IN, a.top/EMU_IN, a.width/EMU_IN, a.height/EMU_IN
            bx, by, bw2, bh2 = b.left/EMU_IN, b.top/EMU_IN, b.width/EMU_IN, b.height/EMU_IN
            ox = max(0, min(ax+aw, bx+bw2) - max(ax, bx))
            oy = max(0, min(ay+ah, by+bh2) - max(ay, by))
            if ox > 0.12 and oy > 0.12:
                issues.append(f"S{si} OVERLAP: '{a.text_frame.text[:18]}' vs '{b.text_frame.text[:18]}' ({ox:.2f}x{oy:.2f}in)")

# placeholder / leftover text scan
pat = re.compile(r"xxx|lorem|TODO|\[insert|undefined|NaN|null", re.I)
for si, slide in enumerate(prs.slides, 1):
    for sh in txt_shapes(slide.shapes):
        t = sh.text_frame.text
        if pat.search(t):
            issues.append(f"S{si} PLACEHOLDER: '{t[:40]}'")

# notes presence + slide count
assert len(prs.slides) == 17, f"slide count {len(prs.slides)} != 17"
missing_notes = [i for i, s in enumerate(prs.slides, 1) if not (s.has_notes_slide and s.notes_slide.notes_text_frame.text.strip())]

print(f"slides: {len(prs.slides)}  |  size: {SLIDE_W:.2f}x{SLIDE_H:.2f}in")
print(f"slides missing speaker notes: {missing_notes if missing_notes else 'none'}")
print(f"issues found: {len(issues)}")
for it in issues:
    print(" -", it)
